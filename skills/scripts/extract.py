"""Stage 1 of deck-to-lecture: pptx -> raw markdown + images + contact sheet.

Everything here is deterministic. No judgement, no LLM. Run this first, then
look at the contact sheet and raw.md, then write map.json and run build.py.

    python extract.py "path/to/Deck.pptx" --work ./work

Outputs into --work:
    raw.md          pptx2md output, unmodified
    img/            every image in the deck, in slide order
    contact.png     numbered grid of all images (look at this ONCE, not N times)
    deck-info.json  slide count, hidden slides, speaker notes, image repeat counts
"""
import argparse, hashlib, json, os, re, shutil, subprocess, sys, tempfile
from collections import Counter

from pptx import Presentation


def slide_is_hidden(slide):
    return slide._element.get("show") == "0"


def strip_hidden(src, dst):
    """Write a copy of the deck with hidden slides removed.

    Hidden slides were deliberately taken out of the live presentation, so they
    must not reach the page. Removing them here (rather than filtering later)
    keeps slide numbering 1:1 with what pptx2md emits.
    """
    prs = Presentation(src)
    hidden = [i for i, s in enumerate(prs.slides, 1) if slide_is_hidden(s)]
    if hidden:
        id_list = prs.slides._sldIdLst
        for el in [list(id_list)[i - 1] for i in hidden]:
            id_list.remove(el)
    prs.save(dst)
    return hidden


def collect_notes(path):
    """{visible_slide_number: notes_text} for slides that survive stripping."""
    prs = Presentation(path)
    out = {}
    for i, s in enumerate(prs.slides, 1):
        if s.has_notes_slide:
            t = s.notes_slide.notes_text_frame.text.strip()
            if t:
                out[i] = t
    return out


def contact_sheet(img_dir, out_path, cols=5, cell=300):
    from PIL import Image, ImageDraw

    def key(p):
        m = re.search(r"_(\d+)\.[^.]+$", os.path.basename(p))
        return int(m.group(1)) if m else 0

    files = sorted(
        (os.path.join(img_dir, f) for f in os.listdir(img_dir)), key=key
    )
    if not files:
        return [], 0
    rows = (len(files) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * cell, rows * cell), (28, 28, 32))
    d = ImageDraw.Draw(sheet)
    hashes = []
    for i, f in enumerate(files):
        h = hashlib.md5(open(f, "rb").read()).hexdigest()
        hashes.append((os.path.basename(f), h[:10], os.path.getsize(f)))
        try:
            im = Image.open(f).convert("RGB")
            im.thumbnail((cell - 16, cell - 40))
            x, y = (i % cols) * cell, (i // cols) * cell
            sheet.paste(im, (x + 8, y + 30))
            d.text((x + 10, y + 8), str(i) + "  " + os.path.basename(f)[-14:],
                   fill=(255, 220, 120))
        except Exception as e:
            print("  !! could not thumbnail", f, e)
    sheet.save(out_path)
    return hashes, len(files)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--work", default="./work")
    a = ap.parse_args()

    work = os.path.abspath(a.work)
    img = os.path.join(work, "img")
    shutil.rmtree(work, ignore_errors=True)
    os.makedirs(img, exist_ok=True)

    total = len(Presentation(a.pptx).slides._sldIdLst)
    # pptx2md names extracted images after the input file, so the stripped copy
    # is called "deck.pptx" -- images then land as deck_0.png, deck_1.jpg, ...
    # Predictable names matter: they are the keys of map.json.
    tmp = os.path.join(tempfile.mkdtemp(prefix="deck2lec-"), "deck.pptx")
    hidden = strip_hidden(a.pptx, tmp)
    notes = collect_notes(tmp)
    visible = len(Presentation(tmp).slides._sldIdLst)

    print("slides: " + str(total) + " total, " + str(len(hidden)) +
          " hidden (dropped: " + str(hidden) + "), " + str(visible) + " kept")
    print("speaker notes on " + str(len(notes)) + " slides")

    raw = os.path.join(work, "raw.md")
    subprocess.run(
        [sys.executable, "-m", "pptx2md", tmp, "-o", raw, "-i", img,
         "--enable-slides"],
        check=True, capture_output=True,
    )

    hashes, n_img = contact_sheet(img, os.path.join(work, "contact.png"))
    repeats = Counter(h for _, h, _ in hashes)
    chrome = [name for name, h, _ in hashes if repeats[h] >= 3]

    info = {
        "source": os.path.basename(a.pptx),
        "slides_total": total,
        "slides_hidden": hidden,
        "slides_visible": visible,
        "notes": {str(k): v for k, v in notes.items()},
        "images": [{"file": n, "hash": h, "bytes": b} for n, h, b in hashes],
        "likely_chrome": sorted(set(chrome)),
    }
    json.dump(info, open(os.path.join(work, "deck-info.json"), "w",
                         encoding="utf-8"), indent=2, ensure_ascii=False)

    mb = sum(b for _, _, b in hashes) / 1048576
    print("images: " + str(n_img) + " (" + format(mb, ".1f") + " MB)")
    if chrome:
        print("likely template chrome (repeats 3+ times): " +
              str(sorted(set(chrome))))
    print()
    print("wrote " + work)
    print("NEXT: look at contact.png, skim raw.md, then write map.json "
          "and run build.py")


if __name__ == "__main__":
    main()
